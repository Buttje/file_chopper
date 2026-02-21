"""
Unit tests for file_chopper.segmenter
"""

from __future__ import annotations

from pathlib import Path
from unittest.mock import patch

import pytest

from file_chopper.segmenter import (
    OUT_OF_SCOPE_FORMATS,
    SUPPORTED_FORMATS,
    SegmentResult,
    SegmentStatus,
    _collect_odf_text,
    _split_text_into_chunks,
    extract_text,
    needs_split,
    segment_document,
    segment_folder,
)

# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture()
def txt_file(tmp_path: Path) -> Path:
    p = tmp_path / "hello.txt"
    p.write_text("Hello, world!", encoding="utf-8")
    return p


@pytest.fixture()
def large_txt_file(tmp_path: Path) -> Path:
    p = tmp_path / "large.txt"
    p.write_text("A" * 200_000, encoding="utf-8")
    return p


@pytest.fixture()
def html_file(tmp_path: Path) -> Path:
    p = tmp_path / "page.html"
    p.write_text(
        "<html><head><title>T</title></head><body>"
        "<script>alert(1)</script>"
        "<style>body{}</style>"
        "<p>Hello HTML</p>"
        "</body></html>",
        encoding="utf-8",
    )
    return p


@pytest.fixture()
def docx_file(tmp_path: Path) -> Path:
    import docx

    doc = docx.Document()
    doc.add_paragraph("First paragraph")
    doc.add_paragraph("Second paragraph")
    p = tmp_path / "test.docx"
    doc.save(str(p))
    return p


@pytest.fixture()
def pptx_file(tmp_path: Path) -> Path:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Slide Title"
    slide.placeholders[1].text = "Slide content"
    p = tmp_path / "test.pptx"
    prs.save(str(p))
    return p


@pytest.fixture()
def xlsx_file(tmp_path: Path) -> Path:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Name"
    ws["B1"] = "Value"
    ws["A2"] = "Alpha"
    ws["B2"] = 42
    p = tmp_path / "test.xlsx"
    wb.save(str(p))
    return p


@pytest.fixture()
def odt_file(tmp_path: Path) -> Path:
    from odf.opendocument import OpenDocumentText
    from odf.text import P

    doc = OpenDocumentText()
    para = P(text="ODT paragraph content")
    doc.text.addElement(para)
    p = tmp_path / "test.odt"
    doc.save(str(p))
    return p


@pytest.fixture()
def rtf_file(tmp_path: Path) -> Path:
    p = tmp_path / "test.rtf"
    # Minimal valid RTF
    p.write_text(
        r"{\rtf1\ansi{\fonttbl\f0\fswiss Helvetica;}\f0\pard Hello RTF\par}",
        encoding="latin-1",
    )
    return p


@pytest.fixture()
def pdf_file(tmp_path: Path) -> Path:
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    p = tmp_path / "test.pdf"
    with p.open("wb") as fh:
        writer.write(fh)
    return p


# ---------------------------------------------------------------------------
# Format sets
# ---------------------------------------------------------------------------


class TestFormatSets:
    def test_doc_in_out_of_scope(self):
        assert ".doc" in OUT_OF_SCOPE_FORMATS

    def test_ppt_in_out_of_scope(self):
        assert ".ppt" in OUT_OF_SCOPE_FORMATS

    def test_xls_in_out_of_scope(self):
        assert ".xls" in OUT_OF_SCOPE_FORMATS

    def test_docx_in_supported(self):
        assert ".docx" in SUPPORTED_FORMATS

    def test_pdf_in_supported(self):
        assert ".pdf" in SUPPORTED_FORMATS

    def test_txt_in_supported(self):
        assert ".txt" in SUPPORTED_FORMATS

    def test_html_in_supported(self):
        assert ".html" in SUPPORTED_FORMATS

    def test_htm_in_supported(self):
        assert ".htm" in SUPPORTED_FORMATS


# ---------------------------------------------------------------------------
# extract_text
# ---------------------------------------------------------------------------


class TestExtractText:
    def test_txt_extraction(self, txt_file: Path):
        assert extract_text(txt_file) == "Hello, world!"

    def test_csv_extraction(self, tmp_path: Path):
        p = tmp_path / "data.csv"
        p.write_text("a,b,c\n1,2,3\n", encoding="utf-8")
        text = extract_text(p)
        assert "a,b,c" in text

    def test_md_extraction(self, tmp_path: Path):
        p = tmp_path / "readme.md"
        p.write_text("# Heading\n\nSome content.", encoding="utf-8")
        text = extract_text(p)
        assert "Heading" in text

    def test_html_extraction(self, html_file: Path):
        text = extract_text(html_file)
        assert "Hello HTML" in text

    def test_html_script_not_included(self, html_file: Path):
        text = extract_text(html_file)
        assert "alert" not in text

    def test_html_style_not_included(self, html_file: Path):
        text = extract_text(html_file)
        assert "body{}" not in text

    def test_htm_extension(self, tmp_path: Path):
        p = tmp_path / "page.htm"
        p.write_text("<html><body><p>HTM content</p></body></html>", encoding="utf-8")
        assert "HTM content" in extract_text(p)

    def test_docx_extraction(self, docx_file: Path):
        text = extract_text(docx_file)
        assert "First paragraph" in text
        assert "Second paragraph" in text

    def test_pptx_extraction(self, pptx_file: Path):
        text = extract_text(pptx_file)
        assert "Slide" in text

    def test_xlsx_extraction(self, xlsx_file: Path):
        text = extract_text(xlsx_file)
        assert "Name" in text
        assert "Alpha" in text

    def test_odt_extraction(self, odt_file: Path):
        text = extract_text(odt_file)
        assert "ODT paragraph" in text

    def test_rtf_extraction(self, rtf_file: Path):
        text = extract_text(rtf_file)
        assert "Hello RTF" in text

    def test_pdf_extraction(self, pdf_file: Path):
        # blank PDF page — just verify it runs without error
        text = extract_text(pdf_file)
        assert isinstance(text, str)

    def test_unsupported_format_raises(self, tmp_path: Path):
        p = tmp_path / "file.xyz"
        p.write_bytes(b"\x00")
        with pytest.raises(ValueError, match="not supported"):
            extract_text(p)

    def test_missing_pypdf_raises_import_error(self, tmp_path: Path):
        p = tmp_path / "file.pdf"
        p.write_bytes(b"%PDF-1.4")
        with patch.dict("sys.modules", {"pypdf": None}), pytest.raises(ImportError, match="pypdf"):
            from file_chopper.segmenter import _extract_pdf_text

            _extract_pdf_text(p)

    def test_missing_docx_raises_import_error(self, tmp_path: Path):
        p = tmp_path / "file.docx"
        p.write_bytes(b"PK")
        with patch.dict("sys.modules", {"docx": None}), pytest.raises(ImportError, match="python-docx"):
            from file_chopper.segmenter import _extract_docx_text

            _extract_docx_text(p)

    def test_missing_pptx_raises_import_error(self, tmp_path: Path):
        p = tmp_path / "file.pptx"
        p.write_bytes(b"PK")
        with patch.dict("sys.modules", {"pptx": None}), pytest.raises(ImportError, match="python-pptx"):
            from file_chopper.segmenter import _extract_pptx_text

            _extract_pptx_text(p)

    def test_missing_openpyxl_raises_import_error(self, tmp_path: Path):
        p = tmp_path / "file.xlsx"
        p.write_bytes(b"PK")
        with patch.dict("sys.modules", {"openpyxl": None}), pytest.raises(ImportError, match="openpyxl"):
            from file_chopper.segmenter import _extract_xlsx_text

            _extract_xlsx_text(p)

    def test_missing_odfpy_raises_import_error(self, tmp_path: Path):
        p = tmp_path / "file.odt"
        p.write_bytes(b"PK")
        odf_modules = {"odf": None, "odf.opendocument": None, "odf.teletype": None, "odf.element": None}
        with patch.dict("sys.modules", odf_modules), pytest.raises(ImportError, match="odfpy"):
            from file_chopper.segmenter import _extract_odf_text

            _extract_odf_text(p)

    def test_missing_striprtf_raises_import_error(self, tmp_path: Path):
        p = tmp_path / "file.rtf"
        p.write_text(r"{\rtf1}", encoding="latin-1")
        striprtf_modules = {"striprtf": None, "striprtf.striprtf": None}
        with patch.dict("sys.modules", striprtf_modules), pytest.raises(ImportError, match="striprtf"):
            from file_chopper.segmenter import _extract_rtf_text

            _extract_rtf_text(p)


# ---------------------------------------------------------------------------
# _collect_odf_text
# ---------------------------------------------------------------------------


class TestCollectOdfText:
    def test_handles_exception_gracefully(self):
        """_collect_odf_text should not raise when teletype fails."""

        class _FakeNode:
            pass

        parts: list = []
        # Should not raise even with an object that has no compatible interface
        _collect_odf_text(_FakeNode(), parts)
        # May or may not append content — just must not raise


# ---------------------------------------------------------------------------
# needs_split
# ---------------------------------------------------------------------------


class TestNeedsSplit:
    def test_size_exceeds_triggers_split(self, txt_file: Path):
        # file is tiny; set max_child_bytes to 1 to force split
        assert needs_split(txt_file, max_child_bytes=1, max_child_text_chars=100_000)

    def test_size_within_limit_no_split(self, txt_file: Path):
        assert not needs_split(
            txt_file, max_child_bytes=1_000_000, max_child_text_chars=100_000
        )

    def test_text_chars_exceeds_triggers_split(self, large_txt_file: Path):
        assert needs_split(
            large_txt_file, max_child_bytes=1_000_000, max_child_text_chars=100
        )

    def test_preextracted_text_used(self, txt_file: Path):
        # precomputed long text should trigger split
        long_text = "X" * 1_000_000
        assert needs_split(
            txt_file,
            max_child_bytes=1_000_000,
            max_child_text_chars=100,
            _text=long_text,
        )

    def test_preextracted_text_short_no_split(self, txt_file: Path):
        assert not needs_split(
            txt_file,
            max_child_bytes=1_000_000,
            max_child_text_chars=100_000,
            _text="short",
        )

    def test_extraction_failure_returns_false(self, tmp_path: Path):
        # Unknown format; extraction raises ValueError → treated as no split
        p = tmp_path / "file.xyz"
        p.write_bytes(b"data")
        assert not needs_split(p, max_child_bytes=1_000_000, max_child_text_chars=100_000)


# ---------------------------------------------------------------------------
# _split_text_into_chunks
# ---------------------------------------------------------------------------


class TestSplitTextIntoChunks:
    def test_single_chunk_when_fits(self):
        result = _split_text_into_chunks("hello", 100)
        assert result == ["hello"]

    def test_splits_evenly(self):
        result = _split_text_into_chunks("ABCDEF", 2)
        assert result == ["AB", "CD", "EF"]

    def test_last_chunk_smaller(self):
        result = _split_text_into_chunks("ABCDE", 2)
        assert result == ["AB", "CD", "E"]

    def test_empty_string_produces_no_chunks(self):
        result = _split_text_into_chunks("", 100)
        assert result == []

    def test_total_reconstructs_original(self):
        text = "X" * 1_000_007
        chunks = _split_text_into_chunks(text, 1_000)
        assert "".join(chunks) == text


# ---------------------------------------------------------------------------
# segment_document
# ---------------------------------------------------------------------------


class TestSegmentDocument:
    # ---- out-of-scope, no split ----

    def test_doc_small_copied(self, tmp_path: Path):
        doc = tmp_path / "small.doc"
        doc.write_bytes(b"DOCBIN" * 10)
        out = tmp_path / "out"
        result = segment_document(doc, out, max_child_bytes=10_000, max_child_text_chars=50_000)
        assert result.status == SegmentStatus.OK
        assert (out / "small.doc").read_bytes() == doc.read_bytes()

    def test_ppt_small_copied(self, tmp_path: Path):
        f = tmp_path / "slides.ppt"
        f.write_bytes(b"X" * 100)
        out = tmp_path / "out"
        result = segment_document(f, out, max_child_bytes=10_000, max_child_text_chars=50_000)
        assert result.status == SegmentStatus.OK

    def test_xls_small_copied(self, tmp_path: Path):
        f = tmp_path / "data.xls"
        f.write_bytes(b"X" * 100)
        out = tmp_path / "out"
        result = segment_document(f, out, max_child_bytes=10_000, max_child_text_chars=50_000)
        assert result.status == SegmentStatus.OK

    def test_doc_large_error(self, tmp_path: Path):
        doc = tmp_path / "large.doc"
        doc.write_bytes(b"X" * 20_000)
        out = tmp_path / "out"
        result = segment_document(doc, out, max_child_bytes=10_000, max_child_text_chars=50_000)
        assert result.status == SegmentStatus.ERROR
        assert result.children == []

    # ---- unknown format (not out-of-scope, not supported) ----

    def test_unknown_format_copied(self, tmp_path: Path):
        f = tmp_path / "data.bin"
        f.write_bytes(b"binary data")
        out = tmp_path / "out"
        result = segment_document(f, out, max_child_bytes=10_000, max_child_text_chars=50_000)
        assert result.status == SegmentStatus.OK
        assert len(result.children) == 1

    # ---- supported formats, no split ----

    def test_txt_small_copied(self, txt_file: Path, tmp_path: Path):
        out = tmp_path / "out"
        result = segment_document(
            txt_file, out, max_child_bytes=10_000, max_child_text_chars=50_000
        )
        assert result.status == SegmentStatus.OK
        assert len(result.children) == 1
        assert result.children[0].suffix == ".txt"

    def test_html_small_copied(self, html_file: Path, tmp_path: Path):
        out = tmp_path / "out"
        result = segment_document(
            html_file, out, max_child_bytes=10_000, max_child_text_chars=50_000
        )
        assert result.status == SegmentStatus.OK
        assert len(result.children) == 1

    # ---- supported formats, split ----

    def test_txt_large_split(self, large_txt_file: Path, tmp_path: Path):
        out = tmp_path / "out"
        result = segment_document(
            large_txt_file, out, max_child_bytes=10 * 1024 * 1024, max_child_text_chars=1_000
        )
        assert result.status == SegmentStatus.OK
        assert len(result.children) > 1
        # All children are .txt
        for child in result.children:
            assert child.suffix == ".txt"

    def test_txt_split_by_file_size(self, tmp_path: Path):
        f = tmp_path / "medium.txt"
        f.write_text("Hello world", encoding="utf-8")
        out = tmp_path / "out"
        # Force split by file size (even though text is short)
        result = segment_document(f, out, max_child_bytes=1, max_child_text_chars=50_000)
        assert result.status == SegmentStatus.OK
        # Should create child files (at least 1)
        assert len(result.children) >= 1

    def test_output_dir_created(self, txt_file: Path, tmp_path: Path):
        out = tmp_path / "deep" / "nested" / "out"
        assert not out.exists()
        segment_document(txt_file, out, max_child_bytes=10_000, max_child_text_chars=50_000)
        assert out.exists()

    # ---- error during text extraction ----

    def test_extraction_exception_returns_error(self, tmp_path: Path):
        """If text extraction raises a non-ImportError exception, result is ERROR."""
        f = tmp_path / "bad.docx"
        f.write_bytes(b"not a real docx")
        out = tmp_path / "out"
        result = segment_document(f, out, max_child_bytes=1, max_child_text_chars=50_000)
        assert result.status == SegmentStatus.ERROR

    def test_missing_dep_returns_missing_dep_status(self, tmp_path: Path):
        """If a required library is missing, result has MISSING_DEP status."""
        f = tmp_path / "file.pdf"
        f.write_bytes(b"%PDF-1.4 minimal")
        out = tmp_path / "out"
        with patch("file_chopper.segmenter.extract_text", side_effect=ImportError("pypdf missing")):
            result = segment_document(f, out, max_child_bytes=1, max_child_text_chars=50_000)
        assert result.status == SegmentStatus.MISSING_DEP

    # ---- ods/odp variants ----

    def test_odp_extraction_via_extract_text(self, tmp_path: Path):
        """segment_document handles .odp (ODF presentation)."""
        from odf.draw import Page
        from odf.opendocument import OpenDocumentPresentation

        doc = OpenDocumentPresentation()
        page = Page(stylename="", masterpagename="")
        doc.presentation.addElement(page)
        p = tmp_path / "slides.odp"
        doc.save(str(p))
        out = tmp_path / "out"
        result = segment_document(p, out, max_child_bytes=10_000, max_child_text_chars=50_000)
        assert result.status == SegmentStatus.OK

    def test_ods_extraction_via_extract_text(self, tmp_path: Path):
        """segment_document handles .ods (ODF spreadsheet)."""
        from odf.opendocument import OpenDocumentSpreadsheet

        doc = OpenDocumentSpreadsheet()
        p = tmp_path / "data.ods"
        doc.save(str(p))
        out = tmp_path / "out"
        result = segment_document(p, out, max_child_bytes=10_000, max_child_text_chars=50_000)
        assert result.status == SegmentStatus.OK


# ---------------------------------------------------------------------------
# segment_folder
# ---------------------------------------------------------------------------


class TestSegmentFolder:
    def test_processes_multiple_files(self, tmp_path: Path):
        input_dir = tmp_path / "input"
        input_dir.mkdir()
        for name in ("a.txt", "b.txt", "c.txt"):
            (input_dir / name).write_text("hello", encoding="utf-8")
        output_dir = tmp_path / "output"

        results = segment_folder(input_dir, output_dir, 10_000, 50_000)
        assert len(results) == 3
        assert all(r.status == SegmentStatus.OK for r in results)

    def test_mirrors_directory_structure(self, tmp_path: Path):
        input_dir = tmp_path / "input"
        sub = input_dir / "sub"
        sub.mkdir(parents=True)
        (sub / "nested.txt").write_text("nested content", encoding="utf-8")
        output_dir = tmp_path / "output"

        segment_folder(input_dir, output_dir, 10_000, 50_000)
        assert (output_dir / "sub" / "nested.txt").exists()

    def test_empty_directory_returns_empty(self, tmp_path: Path):
        input_dir = tmp_path / "empty"
        input_dir.mkdir()
        results = segment_folder(input_dir, tmp_path / "out", 10_000, 50_000)
        assert results == []

    def test_fail_fast_stops_at_first_error(self, tmp_path: Path):
        input_dir = tmp_path / "input"
        input_dir.mkdir()
        # Error file first (alphabetically)
        (input_dir / "aaa.doc").write_bytes(b"X" * 20_000)
        (input_dir / "zzz.txt").write_text("ok", encoding="utf-8")
        results = segment_folder(input_dir, tmp_path / "out", 10_000, 50_000, fail_fast=True)
        assert len(results) == 1
        assert results[0].status == SegmentStatus.ERROR

    def test_no_fail_fast_processes_all(self, tmp_path: Path):
        input_dir = tmp_path / "input"
        input_dir.mkdir()
        (input_dir / "aaa.doc").write_bytes(b"X" * 20_000)
        (input_dir / "zzz.txt").write_text("ok", encoding="utf-8")
        results = segment_folder(input_dir, tmp_path / "out", 10_000, 50_000, fail_fast=False)
        assert len(results) == 2


# ---------------------------------------------------------------------------
# SegmentResult dataclass
# ---------------------------------------------------------------------------


class TestSegmentResult:
    def test_default_children_is_empty_list(self, tmp_path: Path):
        r = SegmentResult(source=tmp_path / "f.txt", status=SegmentStatus.OK)
        assert r.children == []

    def test_default_error_msg_is_empty(self, tmp_path: Path):
        r = SegmentResult(source=tmp_path / "f.txt", status=SegmentStatus.ERROR)
        assert r.error_msg == ""
