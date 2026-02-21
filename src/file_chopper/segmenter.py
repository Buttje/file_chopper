"""
Fess Document Segmenter — core logic.

Splits large source documents into multiple smaller child documents and
prepares an output folder structure suitable for a Fess filesystem crawl.
"""

from __future__ import annotations

import html
import shutil
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path

# ---------------------------------------------------------------------------
# Format classification
# ---------------------------------------------------------------------------

#: Formats that MUST be supported via a pure-Python pipeline.
SUPPORTED_FORMATS: frozenset[str] = frozenset(
    {
        ".pdf",
        ".txt",
        ".csv",
        ".md",
        ".html",
        ".htm",
        ".docx",
        ".pptx",
        ".xlsx",
        ".odt",
        ".odp",
        ".ods",
        ".rtf",
    }
)

#: Formats that are explicitly out-of-scope for v1 (legacy binary Office).
OUT_OF_SCOPE_FORMATS: frozenset[str] = frozenset({".doc", ".ppt", ".xls"})


# ---------------------------------------------------------------------------
# Result types
# ---------------------------------------------------------------------------


class SegmentStatus(str, Enum):
    OK = "ok"
    ERROR = "error"
    MISSING_DEP = "missing_dep"


@dataclass
class SegmentResult:
    """Result of processing a single document."""

    source: Path
    status: SegmentStatus
    children: list[Path] = field(default_factory=list)
    error_msg: str = ""


# ---------------------------------------------------------------------------
# Text extraction helpers
# ---------------------------------------------------------------------------


def _extract_text_plain(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _extract_html_text(path: Path) -> str:
    from html.parser import HTMLParser

    class _Extractor(HTMLParser):
        def __init__(self) -> None:
            super().__init__()
            self._chunks: list[str] = []
            self._skip = False

        def handle_starttag(self, tag: str, attrs) -> None:
            if tag in {"script", "style"}:
                self._skip = True

        def handle_endtag(self, tag: str) -> None:
            if tag in {"script", "style"}:
                self._skip = False

        def handle_data(self, data: str) -> None:
            if not self._skip:
                self._chunks.append(data)

    raw = path.read_text(encoding="utf-8", errors="replace")
    extractor = _Extractor()
    extractor.feed(raw)
    return html.unescape("".join(extractor._chunks))


def _extract_pdf_text(path: Path) -> str:
    try:
        import pypdf  # type: ignore[import-untyped]
    except ImportError as exc:
        raise ImportError(
            "pypdf is required to extract text from PDF files.  "
            "Install it with: pip install pypdf"
        ) from exc

    reader = pypdf.PdfReader(str(path))
    parts = []
    for page in reader.pages:
        text = page.extract_text() or ""
        parts.append(text)
    return "\n".join(parts)


def _extract_docx_text(path: Path) -> str:
    try:
        import docx  # type: ignore[import-untyped]
    except ImportError as exc:
        raise ImportError(
            "python-docx is required to extract text from DOCX files.  "
            "Install it with: pip install python-docx"
        ) from exc

    doc = docx.Document(str(path))
    return "\n".join(para.text for para in doc.paragraphs)


def _extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError as exc:
        raise ImportError(
            "python-pptx is required to extract text from PPTX files.  "
            "Install it with: pip install python-pptx"
        ) from exc

    prs = Presentation(str(path))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    parts.append(para.text)
    return "\n".join(parts)


def _extract_xlsx_text(path: Path) -> str:
    try:
        import openpyxl  # type: ignore[import-untyped]
    except ImportError as exc:
        raise ImportError(
            "openpyxl is required to extract text from XLSX files.  "
            "Install it with: pip install openpyxl"
        ) from exc

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join("" if cell is None else str(cell) for cell in row)
            parts.append(row_text)
    return "\n".join(parts)


def _extract_odf_text(path: Path) -> str:
    try:
        from odf import opendocument  # type: ignore[import-untyped]
    except ImportError as exc:
        raise ImportError(
            "odfpy is required to extract text from ODF files.  "
            "Install it with: pip install odfpy"
        ) from exc

    doc = opendocument.load(str(path))
    parts = []
    for elem in doc.body.childNodes:
        _collect_odf_text(elem, parts)
    return "\n".join(parts)


def _collect_odf_text(node, parts: list[str]) -> None:
    from odf import teletype  # type: ignore[import-untyped]

    try:
        text = teletype.extractText(node)
        if text:
            parts.append(text)
    except Exception:
        pass


def _extract_rtf_text(path: Path) -> str:
    try:
        from striprtf.striprtf import rtf_to_text  # type: ignore[import-untyped]
    except ImportError as exc:
        raise ImportError(
            "striprtf is required to extract text from RTF files.  "
            "Install it with: pip install striprtf"
        ) from exc

    raw = path.read_text(encoding="latin-1", errors="replace")
    return rtf_to_text(raw)


def extract_text(path: Path) -> str:
    """Extract plain text from *path*.

    Parameters
    ----------
    path:
        Path to the document.

    Returns
    -------
    str
        Plain text content of the document.

    Raises
    ------
    ImportError
        When a required third-party library is not installed.
    ValueError
        When the format is not supported for text extraction.
    """
    suffix = path.suffix.lower()
    if suffix in {".txt", ".csv", ".md"}:
        return _extract_text_plain(path)
    if suffix in {".html", ".htm"}:
        return _extract_html_text(path)
    if suffix == ".pdf":
        return _extract_pdf_text(path)
    if suffix == ".docx":
        return _extract_docx_text(path)
    if suffix == ".pptx":
        return _extract_pptx_text(path)
    if suffix == ".xlsx":
        return _extract_xlsx_text(path)
    if suffix in {".odt", ".odp", ".ods"}:
        return _extract_odf_text(path)
    if suffix == ".rtf":
        return _extract_rtf_text(path)
    raise ValueError(
        f"Text extraction is not supported for format '{suffix}'."
    )


# ---------------------------------------------------------------------------
# Split decision
# ---------------------------------------------------------------------------


def needs_split(
    path: Path,
    max_child_bytes: int,
    max_child_text_chars: int,
    _text: str | None = None,
) -> bool:
    """Return ``True`` if *path* needs to be split.

    Parameters
    ----------
    path:
        Path to the document.
    max_child_bytes:
        Maximum allowed file size in bytes before splitting is required.
    max_child_text_chars:
        Maximum allowed extracted text length before splitting is required.
    _text:
        Pre-extracted text (avoids double extraction); used internally.
    """
    if path.stat().st_size > max_child_bytes:
        return True
    if _text is not None:
        return len(_text) > max_child_text_chars
    # Extract text to check; if extraction fails treat as no split needed
    try:
        text = extract_text(path)
    except (ImportError, ValueError, Exception):
        return False
    return len(text) > max_child_text_chars


# ---------------------------------------------------------------------------
# Core segmentation
# ---------------------------------------------------------------------------


def _split_text_into_chunks(text: str, max_chars: int) -> list[str]:
    """Split *text* into chunks of at most *max_chars* characters."""
    chunks = []
    start = 0
    while start < len(text):
        chunks.append(text[start : start + max_chars])
        start += max_chars
    return chunks


def segment_document(
    source: Path,
    output_dir: Path,
    max_child_bytes: int,
    max_child_text_chars: int,
) -> SegmentResult:
    """Process a single document.

    Parameters
    ----------
    source:
        Path to the source document.
    output_dir:
        Directory where output files are written.
    max_child_bytes:
        Maximum allowed file size; triggers splitting when exceeded.
    max_child_text_chars:
        Maximum allowed extracted text length; triggers splitting when exceeded.

    Returns
    -------
    SegmentResult
        Contains the processing status and list of output files.
    """
    source = Path(source)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    suffix = source.suffix.lower()

    # ------------------------------------------------------------------ #
    # Out-of-scope formats
    # ------------------------------------------------------------------ #
    if suffix in OUT_OF_SCOPE_FORMATS:
        size_exceeds = source.stat().st_size > max_child_bytes
        if not size_exceeds:
            # Case A: copy unchanged
            dest = output_dir / source.name
            shutil.copy2(source, dest)
            return SegmentResult(source=source, status=SegmentStatus.OK, children=[dest])
        else:
            # Case B: splitting would be required but format not supported
            error_msg = (
                f"Cannot segment '{source}': "
                f"format '{suffix}' is not supported for conversion in v1.\n"
                f"  File path : {source}\n"
                f"  Reason    : legacy binary Office format cannot be converted with Python-only tools\n"
                f"  Remediation: convert the file manually to PDF or use a newer format (.docx)"
            )
            return SegmentResult(source=source, status=SegmentStatus.ERROR, error_msg=error_msg)

    # ------------------------------------------------------------------ #
    # Supported formats
    # ------------------------------------------------------------------ #
    if suffix not in SUPPORTED_FORMATS:
        # Unknown format — copy unchanged (treated as successfully processed)
        dest = output_dir / source.name
        shutil.copy2(source, dest)
        return SegmentResult(source=source, status=SegmentStatus.OK, children=[dest])

    # Extract text; handle missing dependency
    try:
        text = extract_text(source)
    except ImportError as exc:
        return SegmentResult(
            source=source,
            status=SegmentStatus.MISSING_DEP,
            error_msg=str(exc),
        )
    except Exception as exc:
        return SegmentResult(
            source=source,
            status=SegmentStatus.ERROR,
            error_msg=f"Failed to extract text from '{source}': {exc}",
        )

    split_needed = (
        source.stat().st_size > max_child_bytes or len(text) > max_child_text_chars
    )

    if not split_needed:
        # Copy unchanged
        dest = output_dir / source.name
        shutil.copy2(source, dest)
        return SegmentResult(source=source, status=SegmentStatus.OK, children=[dest])

    # Split into child text files
    chunks = _split_text_into_chunks(text, max_child_text_chars)
    children: list[Path] = []
    stem = source.stem
    for i, chunk in enumerate(chunks, start=1):
        child_name = f"{stem}_part{i:04d}.txt"
        child_path = output_dir / child_name
        child_path.write_text(chunk, encoding="utf-8")
        children.append(child_path)

    return SegmentResult(source=source, status=SegmentStatus.OK, children=children)


# ---------------------------------------------------------------------------
# Folder-level segmentation
# ---------------------------------------------------------------------------


def segment_folder(
    input_dir: Path,
    output_dir: Path,
    max_child_bytes: int,
    max_child_text_chars: int,
    fail_fast: bool = False,
) -> list[SegmentResult]:
    """Recursively segment all documents in *input_dir*.

    The output folder structure mirrors *input_dir*.

    Parameters
    ----------
    input_dir:
        Root input directory.
    output_dir:
        Root output directory.
    max_child_bytes:
        Maximum child file size in bytes.
    max_child_text_chars:
        Maximum child text length in characters.
    fail_fast:
        When ``True``, stop processing after the first error.

    Returns
    -------
    list[SegmentResult]
        One result per processed file.
    """
    input_dir = Path(input_dir)
    output_dir = Path(output_dir)
    results: list[SegmentResult] = []

    for source in sorted(input_dir.rglob("*")):
        if not source.is_file():
            continue
        relative = source.relative_to(input_dir)
        child_output_dir = output_dir / relative.parent
        result = segment_document(
            source=source,
            output_dir=child_output_dir,
            max_child_bytes=max_child_bytes,
            max_child_text_chars=max_child_text_chars,
        )
        results.append(result)
        if fail_fast and result.status != SegmentStatus.OK:
            break

    return results
